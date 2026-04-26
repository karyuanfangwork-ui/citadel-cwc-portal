from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colors
BRAND_BLUE = RGBColor(0, 82, 204)      # #0052CC - Citadel brand color
SKY_BLUE = RGBColor(91, 191, 232)      # #5BBFE8 - Citadel sky blue
STEEL_BLUE = RGBColor(74, 141, 184)    # #4A8DB8 - Citadel steel blue
DARK_BG = RGBColor(26, 35, 126)        # Dark blue for title slides
LIGHT_BG = RGBColor(244, 245, 247)     # Light gray for content
WHITE = RGBColor(255, 255, 255)
CHARCOAL = RGBColor(54, 69, 79)        # #36454F
MUTED = RGBColor(102, 102, 102)

def set_slide_background(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_title_box(slide, text, x, y, w, h, font_size=44, color=WHITE, bold=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box

def add_body_text(slide, text, x, y, w, h, font_size=18, color=CHARCOAL, bullets=False, bold=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    if bullets:
        lines = text.split('\n')
        p.level = 0
        for i, line in enumerate(lines[1:], 1):
            if line.strip():
                new_p = tf.add_paragraph()
                new_p.text = line
                new_p.level = 0
                new_p.font.size = Pt(font_size)
                new_p.font.color.rgb = color
    return box

def add_icon_circle(slide, letter, x, y, size=0.8, bg_color=BRAND_BLUE, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = letter
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_flowchart_step(slide, text, x, y, w, h, color=BRAND_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_arrow(slide, x, y, w, h, color=STEEL_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

print('Creating ITSM Employee Guide Presentation...')
print('')

# SLIDE 1: COVER PAGE
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_title_box(slide, 'ITSM Employee User Guide', Inches(0.5), Inches(2.5), Inches(12), Inches(1.5), font_size=48)
add_body_text(slide, 'Your Complete Guide to Internal Support Services', Inches(0.5), Inches(4), Inches(12), Inches(1), font_size=24, color=SKY_BLUE)
add_body_text(slide, 'Citadel Group Technologies Sdn Bhd', Inches(0.5), Inches(6.5), Inches(12), Inches(0.8), font_size=18, color=WHITE)
add_body_text(slide, 'Launch Date: 2026', Inches(0.5), Inches(7.1), Inches(12), Inches(0.4), font_size=14, color=SKY_BLUE)
icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(0.5), Inches(1.5), Inches(1.5))
icon.fill.solid()
icon.fill.fore_color.rgb = SKY_BLUE
icon.line.fill.background()
print('Slide 1: Cover Page')

# SLIDE 2: WHAT IS ITSM
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)
add_title_box(slide, 'What is ITSM?', Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=40, color=DARK_BG)
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(4), Inches(2))
box1.fill.solid()
box1.fill.fore_color.rgb = BRAND_BLUE
box1.line.fill.background()
tf1 = box1.text_frame
tf1.clear()
p1 = tf1.paragraphs[0]
p1.text = 'One Platform\nfor All Internal\nSupport Requests'
p1.font.size = Pt(18)
p1.font.bold = True
p1.font.color.rgb = WHITE
p1.alignment = PP_ALIGN.CENTER
tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(2), Inches(4), Inches(2))
box2.fill.solid()
box2.fill.fore_color.rgb = STEEL_BLUE
box2.line.fill.background()
tf2 = box2.text_frame
tf2.clear()
p2 = tf2.paragraphs[0]
p2.text = 'Faster Response\nTimes\nFrom Our Teams'
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(2), Inches(4), Inches(2))
box3.fill.solid()
box3.fill.fore_color.rgb = SKY_BLUE
box3.line.fill.background()
tf3 = box3.text_frame
tf3.clear()
p3 = tf3.paragraphs[0]
p3.text = 'Easy Tracking\nof Your\nTicket Status'
p3.font.size = Pt(18)
p3.font.bold = True
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER
tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
add_body_text(slide, 'ITSM = Information Technology Service Management\nA simple way to request help and track progress', Inches(0.5), Inches(4.5), Inches(12), Inches(1), font_size=16, color=MUTED)
print('Slide 2: What is ITSM')

# SLIDE 3: SUPPORT SERVICES COVERED
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)
add_title_box(slide, 'Support Services Covered', Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=40, color=DARK_BG)
it_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(4.1), Inches(5))
it_box.fill.solid()
it_box.fill.fore_color.rgb = WHITE
it_box.line.color.rgb = BRAND_BLUE
it_box.line.width = Pt(3)
tf = it_box.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = 'IT Support'
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = BRAND_BLUE
for item in ['Get IT Help', 'Email Management', 'Report System Problem', 'Request Software Installation', 'Request New Hardware']:
    new_p = tf.add_paragraph()
    new_p.text = '\u2022 ' + item
    new_p.font.size = Pt(16)
    new_p.font.color.rgb = CHARCOAL
    new_p.space_before = Pt(8)
hr_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(1.8), Inches(4.1), Inches(5))
hr_box.fill.solid()
hr_box.fill.fore_color.rgb = WHITE
hr_box.line.color.rgb = STEEL_BLUE
hr_box.line.width = Pt(3)
tf = hr_box.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = 'HR Services'
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = STEEL_BLUE
for item in ['Onboarding Requests', 'Document Requests', 'Policy Inquiries', 'Employee Records']:
    new_p = tf.add_paragraph()
    new_p.text = '\u2022 ' + item
    new_p.font.size = Pt(16)
    new_p.font.color.rgb = CHARCOAL
    new_p.space_before = Pt(8)
fin_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.8), Inches(4.1), Inches(5))
fin_box.fill.solid()
fin_box.fill.fore_color.rgb = WHITE
fin_box.line.color.rgb = SKY_BLUE
fin_box.line.width = Pt(3)
tf = fin_box.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = 'Finance Support'
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = SKY_BLUE
for item in ['Expense Claims', 'Budget Requests', 'Financial Reports']:
    new_p = tf.add_paragraph()
    new_p.text = '\u2022 ' + item
    new_p.font.size = Pt(16)
    new_p.font.color.rgb = CHARCOAL
    new_p.space_before = Pt(8)
print('Slide 3: Support Services Covered')

# SLIDE 4: HOW TO SUBMIT A TICKET
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)
add_title_box(slide, 'How to Submit a Ticket', Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=40, color=DARK_BG)
steps = [
    ('1. Log In', 'Access the CWC portal with your company credentials'),
    ('2. Choose Category', 'Select IT, HR, or Finance based on your need'),
    ('3. Fill the Form', 'Provide clear details about your request'),
    ('4. Set Urgency', 'Choose LOW, MEDIUM, HIGH, or CRITICAL'),
    ('5. Submit', 'Click submit and note your ticket number')
]
for i, (title, desc) in enumerate(steps):
    y_pos = Inches(1.8) + i * Inches(1.1)
    add_icon_circle(slide, str(i+1), Inches(0.5), y_pos, size=0.7, bg_color=BRAND_BLUE)
    add_body_text(slide, title, Inches(1.4), y_pos, Inches(3), Inches(0.5), font_size=18, color=DARK_BG, bold=True)
    add_body_text(slide, desc, Inches(1.4), y_pos + Inches(0.4), Inches(11), Inches(0.5), font_size=14, color=CHARCOAL)
placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9), Inches(1.8), Inches(4), Inches(4.5))
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = WHITE
placeholder.line.color.rgb = MUTED
placeholder.line.dash_style = 4
tf = placeholder.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = '[Screenshot:\nTicket Form]'
p.font.size = Pt(16)
p.font.color.rgb = MUTED
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
print('Slide 4: How to Submit a Ticket')

# SLIDE 5: IT REQUEST FLOWCHART
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)
add_title_box(slide, 'IT Request Flowchart', Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=40, color=DARK_BG)
add_flowchart_step(slide, 'Submit Ticket', Inches(0.5), Inches(1.8), Inches(2.5), Inches(1), color=BRAND_BLUE)
add_arrow(slide, Inches(3.1), Inches(2.05), Inches(0.8), Inches(0.5), color=STEEL_BLUE)
add_flowchart_step(slide, 'IT Team Review', Inches(4), Inches(1.8), Inches(2.5), Inches(1), color=STEEL_BLUE)
add_arrow(slide, Inches(6.6), Inches(2.05), Inches(0.8), Inches(0.5), color=STEEL_BLUE)
add_flowchart_step(slide, 'Assignment to\nSpecialist', Inches(7.5), Inches(1.8), Inches(2.5), Inches(1), color=STEEL_BLUE)
add_arrow(slide, Inches(10.1), Inches(2.05), Inches(0.8), Inches(0.5), color=STEEL_BLUE)
add_flowchart_step(slide, 'Resolution &\nClosure', Inches(11), Inches(1.8), Inches(2), Inches(1), color=SKY_BLUE)
statuses = ['SUBMITTED', 'IN PROGRESS', 'RESOLVED', 'CLOSED']
for i, status in enumerate(statuses):
    x_pos = Inches(0.5) + i * Inches(3.2)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(3.5), Inches(2.8), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = [BRAND_BLUE, STEEL_BLUE, SKY_BLUE, RGBColor(100, 100, 100)][i]
    badge.line.fill.background()
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = status
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.8), Inches(12.5), Inches(2))
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = WHITE
placeholder.line.color.rgb = MUTED
placeholder.line.dash_style = 4
tf = placeholder.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = '[Screenshot: IT Ticket Status Tracking View]'
p.font.size = Pt(16)
p.font.color.rgb = MUTED
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
print('Slide 5: IT Request Flowchart')

# SLIDE 6: HR REQUEST FLOWCHART
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)
add_title_box(slide, 'HR Request Flowchart', Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=40, color=DARK_BG)
add_flowchart_step(slide, 'Submit Request', Inches(0.5), Inches(1.8), Inches(2.5), Inches(1), color=STEEL_BLUE)
add_arrow(slide, Inches(3.1), Inches(2.05), Inches(0.8), Inches(0.5), color=STEEL_BLUE)
add_flowchart_step(slide, 'HR Review', Inches(4), Inches(1.8), Inches(2.5), Inches(1), color=STEEL_BLUE)
add_arrow(slide, Inches(6.6), Inches(2.05), Inches(0.8), Inches(0.5), color=STEEL_BLUE)
add_flowchart_step(slide, 'Document\nPreparation', Inches(7.5), Inches(1.8), Inches(2.5), Inches(1), color=STEEL_BLUE)
add_arrow(slide, Inches(10.1), Inches(2.05), Inches(0.8), Inches(0.5), color=STEEL_BLUE)
add_flowchart_step(slide, 'Delivery &\nCompletion', Inches(11), Inches(1.8), Inches(2), Inches(1), color=SKY_BLUE)
statuses = ['SUBMITTED', 'UNDER REVIEW', 'PROCESSING', 'COMPLETED']
for i, status in enumerate(statuses):
    x_pos = Inches(0.5) + i * Inches(3.2)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(3.5), Inches(2.8), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = [STEEL_BLUE, STEEL_BLUE, SKY_BLUE, RGBColor(100, 100, 100)][i]
    badge.line.fill.background()
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = status
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.8), Inches(12.5), Inches(2))
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = WHITE
placeholder.line.color.rgb = MUTED
placeholder.line.dash_style = 4
tf = placeholder.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = '[Screenshot: HR Request Dashboard]'
p.font.size = Pt(16)
p.font.color.rgb = MUTED
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
print('Slide 6: HR Request Flowchart')

# SLIDE 7: FINANCE REQUEST FLOWCHART
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)
add_title_box(slide, 'Finance Request Flowchart', Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=40, color=DARK_BG)
add_flowchart_step(slide, 'Submit Claim/\nRequest', Inches(0.5), Inches(1.8), Inches(2.5), Inches(1), color=SKY_BLUE)
add_arrow(slide, Inches(3.1), Inches(2.05), Inches(0.8), Inches(0.5), color=STEEL_BLUE)
add_flowchart_step(slide, 'Manager\nApproval', Inches(4), Inches(1.8), Inches(2.5), Inches(1), color=STEEL_BLUE)
add_arrow(slide, Inches(6.6), Inches(2.05), Inches(0.8), Inches(0.5), color=STEEL_BLUE)
add_flowchart_step(slide, 'Finance\nVerification', Inches(7.5), Inches(1.8), Inches(2.5), Inches(1), color=STEEL_BLUE)
add_arrow(slide, Inches(10.1), Inches(2.05), Inches(0.8), Inches(0.5), color=STEEL_BLUE)
add_flowchart_step(slide, 'Payment/\nProcessing', Inches(11), Inches(1.8), Inches(2), Inches(1), color=RGBColor(100, 100, 100))
statuses = ['SUBMITTED', 'PENDING APPROVAL', 'VERIFICATION', 'PROCESSED']
for i, status in enumerate(statuses):
    x_pos = Inches(0.5) + i * Inches(3.2)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(3.5), Inches(2.8), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = [SKY_BLUE, STEEL_BLUE, STEEL_BLUE, RGBColor(100, 100, 100)][i]
    badge.line.fill.background()
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = status
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.8), Inches(12.5), Inches(2))
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = WHITE
placeholder.line.color.rgb = MUTED
placeholder.line.dash_style = 4
tf = placeholder.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = '[Screenshot: Finance Request Status]'
p.font.size = Pt(16)
p.font.color.rgb = MUTED
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
print('Slide 7: Finance Request Flowchart')

# SLIDE 8: HOW TO TRACK TICKET STATUS
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)
add_title_box(slide, 'How to Track Your Ticket Status', Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=40, color=DARK_BG)
methods = [
    ('Dashboard View', 'See all your tickets in one place with color-coded status badges'),
    ('Email Notifications', 'Get automatic updates when your ticket status changes'),
    ('Ticket Details Page', 'Click any ticket to view full history and comments'),
    ('Search & Filter', 'Find tickets by category, date, or status quickly')
]
for i, (title, desc) in enumerate(methods):
    y_pos = Inches(1.8) + i * Inches(1.4)
    add_icon_circle(slide, chr(65+i), Inches(0.5), y_pos, size=0.7, bg_color=BRAND_BLUE)
    add_body_text(slide, title, Inches(1.4), y_pos, Inches(4), Inches(0.5), font_size=18, color=DARK_BG, bold=True)
    add_body_text(slide, desc, Inches(1.4), y_pos + Inches(0.4), Inches(11), Inches(0.5), font_size=14, color=CHARCOAL)
add_body_text(slide, 'Status Color Guide:', Inches(8.5), Inches(1.8), Inches(4.5), Inches(0.4), font_size=16, color=DARK_BG, bold=True)
colors = [
    (BRAND_BLUE, 'SUBMITTED - Received'),
    (STEEL_BLUE, 'IN PROGRESS - Being worked on'),
    (SKY_BLUE, 'PENDING - Awaiting action'),
    (RGBColor(100, 100, 100), 'COMPLETED/CLOSED - Finished')
]
for i, (color, label) in enumerate(colors):
    y_pos = Inches(2.4) + i * Inches(0.7)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), y_pos, Inches(0.8), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_body_text(slide, label, Inches(9.5), y_pos, Inches(3.5), Inches(0.5), font_size=14, color=CHARCOAL)
print('Slide 8: How to Track Ticket Status')

# SLIDE 9: WHERE TO GET HELP
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, LIGHT_BG)
add_title_box(slide, 'Where to Get Help', Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=40, color=DARK_BG)
teams = [
    ('IT Team', 'it-support@citadel.com', 'For all technology issues, software requests, and system problems', BRAND_BLUE),
    ('HR Team', 'hr-services@citadel.com', 'For onboarding, documents, policies, and employee records', STEEL_BLUE),
    ('Finance Team', 'finance-helpdesk@citadel.com', 'For expense claims, budget requests, and financial reports', SKY_BLUE)
]
for i, (name, email, desc, color) in enumerate(teams):
    y_pos = Inches(1.8) + i * Inches(2)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y_pos, Inches(12.5), Inches(1.7))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(3)
    add_body_text(slide, name, Inches(1), y_pos + Inches(0.2), Inches(4), Inches(0.5), font_size=22, color=color, bold=True)
    add_body_text(slide, email, Inches(1), y_pos + Inches(0.6), Inches(4), Inches(0.4), font_size=14, color=MUTED)
    add_body_text(slide, desc, Inches(1), y_pos + Inches(1), Inches(11), Inches(0.5), font_size=14, color=CHARCOAL)
    add_icon_circle(slide, name[0], Inches(11.5), y_pos + Inches(0.5), size=1, bg_color=color)
print('Slide 9: Where to Get Help')

# SLIDE 10: QUICK REFERENCE & TIPS
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_title_box(slide, 'Quick Reference & Tips', Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=40, color=WHITE)
tips = [
    '\u2713 Be specific in your request description - more details = faster resolution',
    '\u2713 Set the correct urgency level - CRITICAL is for system outages only',
    '\u2713 Save your ticket number for easy reference in follow-ups',
    '\u2713 Check your email for status updates and requests for more information',
    '\u2713 Close tickets once your issue is resolved to help our teams track metrics'
]
for i, tip in enumerate(tips):
    add_body_text(slide, tip, Inches(0.5), Inches(1.8) + i * Inches(1), Inches(12), Inches(0.8), font_size=18, color=WHITE)
add_body_text(slide, 'Thank you for using CWC ITSM Portal!', Inches(0.5), Inches(6.8), Inches(12), Inches(0.6), font_size=20, color=SKY_BLUE)
print('Slide 10: Quick Reference & Tips')

# Save
prs.save('/Users/fangkaryuan/cwc2.0/citadel-cwc-portal/itsm-employee-guide.pptx')
print('')
print('=' * 50)
print('PRESENTATION CREATED SUCCESSFULLY')
print('=' * 50)
print(f'File: itsm-employee-guide.pptx')
print(f'Total slides: {len(prs.slides)}')
print('')
print('Slide Summary:')
print('  1. Cover Page')
print('  2. What is ITSM')
print('  3. Support Services Covered')
print('  4. How to Submit a Ticket')
print('  5. IT Request Flowchart')
print('  6. HR Request Flowchart')
print('  7. Finance Request Flowchart')
print('  8. How to Track Ticket Status')
print('  9. Where to Get Help')
print(' 10. Quick Reference & Tips')
